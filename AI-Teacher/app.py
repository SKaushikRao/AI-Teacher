from flask import Flask, request, jsonify, render_template
from pptx import Presentation
import os
import fitz  # PyMuPDF for PDFs
from transformers import pipeline

app = Flask(__name__, template_folder="templates")

UPLOAD_FOLDER = "uploads"
ALLOWED_EXTENSIONS = {"pdf", "pptx"}
app.config["UPLOAD_FOLDER"] = UPLOAD_FOLDER

os.makedirs(UPLOAD_FOLDER, exist_ok=True)

# Load BART summarization model
explainer = pipeline("summarization", model="facebook/bart-large-cnn")

def allowed_file(filename):
    return "." in filename and filename.rsplit(".", 1)[1].lower() in ALLOWED_EXTENSIONS

def extract_text_from_pdf(filepath):
    """Extract text from a PDF file."""
    doc = fitz.open(filepath)
    text = "\n".join([page.get_text("text") for page in doc])
    return text

def extract_text_from_pptx(filepath):
    """Extract text from a PPTX file."""
    prs = Presentation(filepath)
    text = "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
    return text

def explain_text(text):
    """Generate an explanation like a teacher would."""
    try:
        max_input_length = 1024  
        text = text[:max_input_length] if len(text) > max_input_length else text

        explanation = explainer(text, max_length=400, min_length=100, do_sample=False)
        return explanation[0]['summary_text']
    except Exception as e:
        print(f"Explanation Error: {e}")
        return "Error in generating explanation."

@app.route('/')
def home():
    return render_template("index.html")

@app.route('/upload', methods=['POST'])
def upload_file():
    """Handle file uploads and generate explanation."""
    if 'file' not in request.files:
        return jsonify({"error": "No file part"}), 400

    file = request.files['file']
    
    if file.filename == '':
        return jsonify({"error": "No selected file"}), 400

    if file and allowed_file(file.filename):
        filepath = os.path.join(app.config["UPLOAD_FOLDER"], file.filename)
        file.save(filepath)

        extracted_text = extract_text_from_pdf(filepath) if file.filename.endswith(".pdf") else extract_text_from_pptx(filepath)

        explained_text = explain_text(extracted_text)

        return jsonify({
            "message": "File uploaded successfully",
            "filename": file.filename,
            "extracted_text": extracted_text[:500],  # Limit for debugging
            "explained_text": explained_text
        }), 200

    return jsonify({"error": "Invalid file format"}), 400

if __name__ == "__main__":
    app.run(debug=True)
